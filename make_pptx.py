"""Generate PulseCare 3-slide pitch deck as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x6B, 0x9E, 0x78)
DARK_GREEN  = RGBColor(0x3D, 0x6B, 0x47)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1C, 0x1C, 0x1E)
GRAY        = RGBColor(0x6B, 0x70, 0x80)
LIGHT_GREEN = RGBColor(0xEA, 0xF4, 0xEE)
ACCENT      = RGBColor(0xC9, 0x78, 0x5A)

W = Inches(13.33)
H = Inches(7.5)

FOOTER_TEXT = "IEORE4576  TOPICS IN OPERATIONS RESEARCH"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = "Helvetica Neue"
    return txb


def footer(slide):
    add_text(slide, FOOTER_TEXT,
             x=Inches(6.5), y=Inches(7.1), w=Inches(6.6), h=Inches(0.35),
             size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def green_bar(slide):
    add_rect(slide, 0, 0, W, Inches(0.08), fill=GREEN)


def section_header(slide, title, subtitle=None):
    green_bar(slide)
    add_text(slide, title,
             x=Inches(0.6), y=Inches(0.22), w=Inches(11), h=Inches(0.65),
             size=30, bold=True, color=DARK_GREEN)
    if subtitle:
        add_text(slide, subtitle,
                 x=Inches(0.6), y=Inches(0.85), w=Inches(12), h=Inches(0.38),
                 size=13, color=GRAY, italic=True)
    rule_y = Inches(1.22) if subtitle else Inches(0.9)
    add_rect(slide, Inches(0.6), rule_y, Inches(12.1), Pt(1.5), fill=LIGHT_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HEADER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Dark green background + lighter right panel
add_rect(s, 0, 0, W, H, fill=DARK_GREEN)
add_rect(s, Inches(8.2), 0, Inches(5.13), H, fill=GREEN)

# Pulse icon circle
add_rect(s, Inches(9.6), Inches(1.8), Inches(2.13), Inches(2.13), fill=DARK_GREEN)
add_text(s, "💚", Inches(9.9), Inches(1.85), Inches(2), Inches(2),
         size=80, align=PP_ALIGN.CENTER, color=WHITE)

# Tag chip
add_rect(s, Inches(0.7), Inches(1.65), Inches(3.6), Inches(0.42),
         fill=RGBColor(0x2A, 0x4D, 0x35))
add_text(s, "AI-Powered Eldercare Platform",
         Inches(0.72), Inches(1.65), Inches(3.56), Inches(0.42),
         size=12, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# Title + subtitle
add_text(s, "PulseCare",
         Inches(0.7), Inches(2.15), Inches(7.2), Inches(1.1),
         size=64, bold=True, color=WHITE)
add_text(s, "AI-Powered Senior Health Monitoring",
         Inches(0.7), Inches(3.2), Inches(7.2), Inches(0.65),
         size=24, color=RGBColor(0xB8, 0xD8, 0xC0))

# Divider
add_rect(s, Inches(0.7), Inches(3.95), Inches(2.4), Pt(3),
         fill=RGBColor(0x6B, 0xC9, 0x9C))

# Byline + tagline
add_text(s, "February Jiang  ·  yj2918",
         Inches(0.7), Inches(4.12), Inches(6), Inches(0.4),
         size=14, color=RGBColor(0xA0, 0xC4, 0xA8))
add_text(s, '"Your body has a story. Pulse remembers it."',
         Inches(0.7), Inches(5.0), Inches(7), Inches(0.5),
         size=14, italic=True, color=RGBColor(0x88, 0xB8, 0x98))

footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVATION & POSITIONING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, "Motivation & Positioning",
               "72M seniors hide pain from families — crises arrive without warning.")

# Personal story card
add_rect(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(1.3),
         fill=RGBColor(0xFD, 0xF5, 0xF2))
add_rect(s, Inches(0.6), Inches(1.35), Pt(5), Inches(1.3), fill=ACCENT)
add_text(s, "Five years ago, I lost my grandmother.",
         Inches(0.85), Inches(1.42), Inches(11.5), Inches(0.4),
         size=15, bold=True, color=DARK)
add_text(s, "She hid her physical pain to avoid worrying our family. "
            "That silence — and the fall that followed — is why PulseCare exists.",
         Inches(0.85), Inches(1.82), Inches(11.5), Inches(0.7),
         size=13, color=GRAY)

# Three problem boxes
problems = [
    ("😶", "Seniors Under-Report",
     "72M independent US seniors systematically hide symptoms to avoid burdening family.\n"
     "Disclosure only happens when pain is severe enough to admit — often too late."),
    ("📞", "Signals Vanish Between Calls",
     "Weekly phone calls miss slow-moving patterns — 14 days of sleep debt, recurring dizziness —"
     " that build silently into a crisis."),
    ("🎯", "PulseCare Closes the Gap",
     "Daily conversational check-ins extract structured health signals invisibly.\n"
     "Caregivers see synthesis, not noise — ranked alerts + predictive 48-hour insight."),
]
box_w = Inches(3.8)
for i, (emoji, title, body) in enumerate(problems):
    bx = Inches(0.6) + i * Inches(4.05)
    by = Inches(2.82)
    add_rect(s, bx, by, box_w, Inches(3.85),
             fill=LIGHT_GREEN if i < 2 else RGBColor(0xE8, 0xF5, 0xEC),
             line=GREEN)
    add_rect(s, bx, by, box_w, Inches(0.1), fill=GREEN if i == 2 else ACCENT)
    add_text(s, emoji, bx, by + Inches(0.18), box_w, Inches(0.65),
             size=30, align=PP_ALIGN.CENTER)
    add_text(s, title, bx + Inches(0.18), by + Inches(0.9),
             box_w - Inches(0.36), Inches(0.5),
             size=13, bold=True, color=DARK_GREEN if i == 2 else ACCENT)
    add_text(s, body, bx + Inches(0.18), by + Inches(1.45),
             box_w - Inches(0.36), Inches(2.2),
             size=11.5, color=GRAY)

# Competitive positioning strip
add_rect(s, Inches(0.6), Inches(6.78), Inches(12.1), Inches(0.55),
         fill=DARK_GREEN)
add_text(s, "vs. Life Alert (hardware stigma)  ·  Apple Watch (raw data, no context)  ·  "
            "ElliQ ($40/mo, no predictive layer)   →   PulseCare: $29/mo · software-only · predictive",
         Inches(0.75), Inches(6.83), Inches(12.0), Inches(0.45),
         size=11, color=WHITE, align=PP_ALIGN.CENTER)

footer(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS MODEL & TECHNICAL BUILDING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, "Business Model & Technical Building",
               "High-margin SaaS powered by conditional multi-agent AI — cost scales sub-linearly with users.")

# ── LEFT: Unit Economics ──────────────────────────────────────────────────────
add_rect(s, Inches(0.6), Inches(1.42), Inches(5.5), Inches(5.7), fill=DARK_GREEN)

add_text(s, "Unit Economics", Inches(0.75), Inches(1.52), Inches(5.2), Inches(0.45),
         size=14, bold=True, color=LIGHT_GREEN)

# Big margin number
add_text(s, "98.8%", Inches(0.6), Inches(1.92), Inches(5.5), Inches(1.25),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Gross Margin", Inches(0.6), Inches(3.12), Inches(5.5), Inches(0.4),
         size=15, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# Cost rows
econ = [
    ("$29.00 /mo",  "Subscription per caregiver"),
    ("~$0.06",      "LLM cost  (Gemini 2.0 Flash, conditional routing)"),
    ("~$0.36",      "Total COGS  (infra + compute)"),
    ("$28.64",      "Contribution margin per user"),
]
ey = Inches(3.65)
for val, lbl in econ:
    add_rect(s, Inches(0.65), ey, Inches(5.4), Inches(0.52),
             fill=RGBColor(0x2A, 0x4D, 0x35))
    add_text(s, val, Inches(0.75), ey + Inches(0.06), Inches(1.9), Inches(0.42),
             size=13, bold=True, color=WHITE)
    add_text(s, lbl, Inches(2.65), ey + Inches(0.09), Inches(3.3), Inches(0.38),
             size=11, color=LIGHT_GREEN)
    ey += Inches(0.58)

add_text(s, "🔒 Moat: after 60 days the senior's longitudinal baseline\nis the product — irreplaceable clinical memory.",
         Inches(0.75), Inches(6.6), Inches(5.2), Inches(0.55),
         size=10.5, italic=True, color=RGBColor(0xA0, 0xC4, 0xA8))

# ── RIGHT: Technical Architecture ────────────────────────────────────────────
tech = [
    ("🔀", "Conditional Multi-Agent Routing",
     "LangGraph: Companion Agent handles every check-in cheaply (~$0.001). "
     "AnalystAgent fires only when severity ≥ 5 or symptoms are detected (~$0.015). "
     "A 'I'm fine' costs nothing meaningful."),
    ("🔬", "Cross-Modal Synthesis",
     "AnalystAgent fuses 60-day conversation history with Apple Watch biometrics "
     "(HR, SpO₂, HRV, sleep stages). Divergence between reported status and vitals "
     "is itself the alert — e.g. 'I feel fine' + SpO₂ 93%."),
    ("📚", "Hybrid RAG  (FAISS + BM25)",
     "Responses grounded in curated geriatric knowledge via Reciprocal Rank Fusion. "
     "Eliminates medical hallucinations and substantiates the $29/mo clinical trust premium."),
    ("🛠", "Forced Tool Calling",
     "log_health_entry is force-called before every warm reply — guaranteeing qualitative "
     "conversation converts into structured DB rows powering every downstream alert."),
]
bw = Inches(6.15)
bh = Inches(1.28)
tx = Inches(6.6)
ty = Inches(1.42)
for i, (emoji, title, body) in enumerate(tech):
    by = ty + i * (bh + Inches(0.08))
    add_rect(s, tx, by, bw, bh, fill=LIGHT_GREEN)
    add_rect(s, tx, by, Pt(5), bh, fill=GREEN)
    add_text(s, f"{emoji}  {title}", tx + Inches(0.2), by + Inches(0.1),
             bw - Inches(0.3), Inches(0.38), size=12, bold=True, color=DARK_GREEN)
    add_text(s, body, tx + Inches(0.2), by + Inches(0.48),
             bw - Inches(0.3), bh - Inches(0.55), size=10.5, color=GRAY)

footer(s)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "/Users/february/Desktop/pulse/PulseCare_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved → {out}")
